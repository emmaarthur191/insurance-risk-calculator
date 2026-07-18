import os
import sys
from pathlib import Path
from lxml import etree

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("ERROR: python-pptx is not installed.")
    print("Please run: pip install python-pptx")
    sys.exit(1)

# ── Morph Transition Helper ──────────────────────────────────────────────────
_P_NS  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
_P159  = 'http://schemas.microsoft.com/office/powerpoint/2015/09/main'

def set_morph_transition(slide, speed='slow'):
    """Inject PowerPoint Morph transition XML into a slide.
    Morph matches shapes by their !! prefixed names across consecutive slides,
    animating position, size, colour and text changes per-card."""
    sld = slide._element
    # Remove any existing transition elements
    for t in list(sld):
        if t.tag.endswith('}transition'):
            sld.remove(t)
    trans = etree.SubElement(sld, f'{{{_P_NS}}}transition',
                             attrib={'spd': speed, 'advClick': '1'})
    etree.SubElement(trans, f'{{{_P159}}}morph',
                     attrib={'option': 'byObject'})

# Theme Colors
CHARCOAL = RGBColor(17, 17, 22)      # #111116 Background
CRIMSON = RGBColor(220, 20, 60)      # #DC143C Accent Red
WHITE = RGBColor(255, 255, 255)      # Title & Primary Text
GRAY = RGBColor(165, 170, 181)       # #A5AAB5 Secondary Text
DARK_GRAY = RGBColor(28, 28, 35)     # #1C1C23 Card fill
SILVER = RGBColor(90, 90, 101)       # #5A5A65 Inactive Card Border

# Slide Dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Paths
LOGO_PATH = Path("landing/logo.png")
if not LOGO_PATH.exists():
    LOGO_PATH = Path(r"C:\Users\snype\.gemini\antigravity-ide\brain\8bb506a0-6824-436d-87c4-355a6c8da02b\media__1784331268747.png")
if not LOGO_PATH.exists():
    LOGO_PATH = Path(r"C:\Users\snype\Downloads\insurance_app\landing\logo.png")

OUTPUT_PATH = Path("board_presentation_executive.pptx")

# Initialize Presentation
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Card positions (consistent across all slides to enable Morph transition)
LEFT_MARGIN = Inches(0.75)
RIGHT_MARGIN = Inches(6.98)
CARD_TOP = Inches(1.8)
CARD_WIDTH = Inches(5.6)
CARD_HEIGHT = Inches(4.8)

def set_slide_background(slide):
    """Fills the slide background with the premium dark charcoal theme color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = CHARCOAL

def add_header(slide, title_text, category_text="THE PRUDENTIAL RISK STORY"):
    """Adds a standard header, category tag, APEX TEAM group text, and the Prudential logo."""
    # Top Category Tag
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(8), Inches(0.4))
    txBox.name = "!!category_tag"
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = category_text.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CRIMSON
    p.font.name = "Segoe UI"
    
    # Title
    txBox_title = slide.shapes.add_textbox(Inches(0.75), Inches(0.7), Inches(9), Inches(0.8))
    txBox_title.name = "!!slide_title"
    tf_title = txBox_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    p_title.font.name = "Segoe UI"

    # Logo in the top-right corner
    if LOGO_PATH.exists():
        logo_shape = slide.shapes.add_picture(str(LOGO_PATH), Inches(11.8), Inches(0.3), height=Inches(0.75))
        logo_shape.name = "!!logo_img"
    
    # APEX TEAM Group Text
    txBox_group = slide.shapes.add_textbox(Inches(10.0), Inches(0.5), Inches(1.8), Inches(0.4))
    txBox_group.name = "!!group_name"
    tf_group = txBox_group.text_frame
    p_group = tf_group.paragraphs[0]
    p_group.text = "APEX TEAM"
    p_group.alignment = PP_ALIGN.RIGHT
    p_group.font.size = Pt(11)
    p_group.font.bold = True
    p_group.font.color.rgb = WHITE
    p_group.font.name = "Segoe UI"

def create_card_shape(slide, left, top, width, height, border_color, name):
    """Draws a card shape container with double exclamation morph name."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_GRAY
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    shape.name = name
    return shape

def create_card_text(slide, left, top, width, height, title_text, bullets, name, title_color=CRIMSON, text_color=WHITE):
    """Creates a text frame inside a card container with double exclamation morph name."""
    txBox = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), height - Inches(0.4))
    txBox.name = name
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = "Segoe UI"
    
    for bullet in bullets:
        p2 = tf.add_paragraph()
        p2.text = "• " + bullet
        p2.font.size = Pt(22)
        p2.font.color.rgb = text_color
        p2.font.name = "Segoe UI"
        p2.space_after = Pt(8)
    return txBox

# ==============================================================================
# DEFINING THE PRESENTATION DATA & CONTENT STRUCTURE (6 SLIDES METRIC TRACK)
# ==============================================================================
slides_content = [
    # Padding for 1-based indexing
    None,
    None,
    
    # ══════════════════════════════════════════════════════════════════════
    # CHAPTER 1: "We started by asking 5 hard questions..."
    # ══════════════════════════════════════════════════════════════════════
    {
        "cat": "CHAPTER 1  •  THE INVESTIGATION",
        "title": "Two Portfolios, One Uncomfortable Truth",
        "left_title": "RETAIL: 50,000 INDIVIDUAL POLICYHOLDERS",
        "left_bullets": [
            "We pulled 950,000 rows from the retail database — but discovered every person was copied 19 times. The real book is 50,000 clients.",
            "We assumed smokers were our biggest cost. They're not — they claim at the same rate as everyone else. Age is the real driver.",
            "Three columns — BMI, tenure, marital status — turned out to be computer-generated noise, not real data."
        ],
        "left_color": CRIMSON,
        "left_text_color": WHITE,
        "right_title": "CORPORATE: GROUP EMPLOYEES WITH REAL CLAIMS",
        "right_bullets": [
            "The corporate dataset tracks actual employees across multiple companies, with real claims filed against their policies.",
            "Unlike retail, corporate has verifiable claims history — frequency, severity, and amounts we can audit.",
            "But 24 corporate claims were approved in under 3 days — some with no proper review at all."
        ],
        "right_color": CRIMSON,
        "right_text_color": WHITE
    },

    # ══════════════════════════════════════════════════════════════════════
    # CHAPTER 2: "The numbers revealed a crisis hiding in plain sight..."
    # ══════════════════════════════════════════════════════════════════════
    {
        "cat": "CHAPTER 2  •  THE CRISIS WE FOUND",
        "title": "Both Books Are Bleeding — But for Different Reasons",
        "left_title": "RETAIL: FLAT PRICING IS KILLING US",
        "left_bullets": [
            "Every retail client pays the same premium — whether they're 25 and healthy or 60 with chronic conditions.",
            "Healthy clients are subsidising high-risk ones without knowing it. That's not pricing — that's guessing.",
            "If nothing changes, the retail book alone is heading toward a GHS 596M pricing gap."
        ],
        "left_color": CRIMSON,
        "left_text_color": WHITE,
        "right_title": "CORPORATE: CLAIMS OUTPACE COLLECTIONS 6-TO-1",
        "right_bullets": [
            "On corporate accounts, we collected GHS 7.3M in premiums — but paid out GHS 47.2M in claims.",
            "That's a GHS 39.9M shortfall on corporate business alone. For every GHS 1.00 in, GHS 6.47 goes out.",
            "Combined across both books, the total exposure is GHS 608M — the 'illusion' in our balance sheet."
        ],
        "right_color": SILVER,
        "right_text_color": GRAY
    },

    # ══════════════════════════════════════════════════════════════════════
    # CHAPTER 3: "Then we looked at the data itself..."
    # ══════════════════════════════════════════════════════════════════════
    {
        "cat": "CHAPTER 3  •  WHAT WE UNCOVERED",
        "title": "The Retail Data Had No Claims — So We Built a Way to Predict Them",
        "left_title": "RETAIL: WHY WE NEEDED TO ESTIMATE LOSSES",
        "left_bullets": [
            "The retail database had no claims history at all — no record of who claimed, how often, or how much. So we couldn't just look at the past.",
            "By using a Generalized Linear Model (GLM) framework — the insurance industry standard — we built two models: a Poisson model that estimates how often a client will claim, and a Gamma model that estimates how much each claim will cost.",
            "These GLMs use only what we know at enrolment — age, income, BMI, dependents — to project each client's expected loss."
        ],
        "left_color": CRIMSON,
        "left_text_color": WHITE,
        "right_title": "WE ALSO CHECKED WHICH FACTORS ACTUALLY MATTER",
        "right_bullets": [
            "Before building anything, we ran a correlation analysis and checked for multicollinearity using VIF (Variance Inflation Factor) — this tells us if two factors are saying the same thing.",
            "We found that income and premium move together, and BMI and obesity overlap. Using both would confuse the model, so we kept only the factors that add independent information.",
            "The corporate workbook gave us real claims across 5 linked sheets — that’s where we applied XGBoost, a machine learning classifier, for experience-based risk rating."
        ],
        "right_color": CRIMSON,
        "right_text_color": WHITE
    },

    # ══════════════════════════════════════════════════════════════════════
    # CHAPTER 4: "The problem wasn't just pricing — it was everywhere..."
    # ══════════════════════════════════════════════════════════════════════
    {
        "cat": "CHAPTER 4  •  THE DEEPER PROBLEMS",
        "title": "It Wasn't Just Pricing — Branches and Incentives Are Broken Too",
        "left_title": "SOME BRANCHES ARE BLEEDING MONEY",
        "left_bullets": [
            "Sunyani alone is GHS 60.2M in the red. Koforidua is close behind at GHS 59.8M.",
            "Claims processing is inconsistent — some offices take 12 days, others over 30.",
            "There's no standard process. How fast a claim gets paid depends on which office handles it."
        ],
        "left_color": CRIMSON,
        "left_text_color": WHITE,
        "right_title": "OUR SALES INCENTIVES REWARD THE WRONG BEHAVIOUR",
        "right_bullets": [
            "Agents earn commissions based on how many policies they sell — not how good the clients are.",
            "One agent's clients cost us GHS 23.4M in claims — but only brought in GHS 1.9M in premiums.",
            "The more our agents sell, the more money the company loses. The incentives are backwards."
        ],
        "right_color": SILVER,
        "right_text_color": GRAY
    },

    # ══════════════════════════════════════════════════════════════════════
    # CHAPTER 5: "Here's how we fix it — starting with pricing..."
    # ══════════════════════════════════════════════════════════════════════
    {
        "cat": "CHAPTER 5  •  THE FIX (PART 1)",
        "title": "Different Books Need Different Cures",
        "left_title": "RETAIL: ESTIMATE FIRST, THEN SCORE",
        "left_bullets": [
            "1. Since retail has no claims data, we use Poisson and Gamma GLMs to estimate each client's expected cost — one model for frequency, one for severity.",
            "2. We then score every applicant using a points-based actuarial scorecard — placing them into Low, Medium, High, or Critical risk tiers.",
            "3. Raise retail premiums by 15–35% based on their tier — that’s what closes the GHS 596M gap."
        ],
        "left_color": CRIMSON,
        "left_text_color": WHITE,
        "right_title": "CORPORATE: ML-POWERED EXPERIENCE RATING",
        "right_bullets": [
            "4. Corporate renewals have real claims data — so we deploy a 4-tier XGBoost model with 90.88% accuracy.",
            "5. The retail model reaches 95.73% accuracy — both classify into Low, Medium, High, and Critical tiers.",
            "6. This hybrid approach (scorecard for new, ML for renewals) saves 3 days per decision at GHS 1 per application."
        ],
        "right_color": CRIMSON,
        "right_text_color": WHITE
    },

    # ══════════════════════════════════════════════════════════════════════
    # CHAPTER 6: "And here's how we protect the company going forward..."
    # ══════════════════════════════════════════════════════════════════════
    {
        "cat": "CHAPTER 6  •  THE FIX (PART 2)",
        "title": "Step Two: Fix Incentives, Catch Fraud, and Build a Safety Net",
        "left_title": "REWARD AGENTS FOR QUALITY — NOT JUST VOLUME",
        "left_bullets": [
            "1. Link 30% of agent commissions to how their clients actually perform over time.",
            "2. No more same-day payouts — hold every claim for 7 days so it can be properly reviewed.",
            "3. Build an auto-flag system that catches suspicious claims before they're paid out."
        ],
        "left_color": CRIMSON,
        "left_text_color": WHITE,
        "right_title": "BUILD A FINANCIAL SAFETY NET",
        "right_bullets": [
            "4. Set aside GHS 544.7M as a reserve fund — this protects us when claims spike unexpectedly.",
            "5. This level of reserves meets international insurance safety standards.",
            "6. Any surplus cash should be invested in safe bonds to grow the fund over time."
        ],
        "right_color": CRIMSON,
        "right_text_color": WHITE
    }
]

# ==============================================================================
# MAIN RENDERER (STAGGERED CARD MORPH ENGINE — TikTok-style reveal)
# Each content slide becomes TWO physical slides:
#   Sub-slide A: LEFT card in position, RIGHT card off-screen (below viewport)
#   Sub-slide B: BOTH cards in position — RIGHT card morphs up into view
# ==============================================================================
OFF_SCREEN_TOP = Inches(8.5)  # Below the visible 7.5" slide height

def build_title_slide():
    """Slide 1: Title page with both card containers staged off-screen."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    if LOGO_PATH.exists():
        logo_title = slide.shapes.add_picture(str(LOGO_PATH), Inches(5.4), Inches(1.2), height=Inches(1.5))
        logo_title.name = "!!logo_img"

    tx_title = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.33), Inches(2.0))
    tx_title.name = "!!slide_title"
    tf1 = tx_title.text_frame
    tf1.word_wrap = True
    p_main = tf1.paragraphs[0]
    p_main.text = "The GHS 608 Million Illusion"
    p_main.alignment = PP_ALIGN.CENTER
    p_main.font.size = Pt(46)
    p_main.font.bold = True
    p_main.font.color.rgb = CRIMSON
    p_main.font.name = "Segoe UI"

    p_sub = tf1.add_paragraph()
    p_sub.text = "Something didn't add up in our books. So we investigated."
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = WHITE
    p_sub.font.name = "Segoe UI"

    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(5.0), Inches(4.33), Inches(0.05))
    accent_line.name = "!!accent_line"
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = CRIMSON
    accent_line.line.fill.background()

    tx_group1 = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.33), Inches(0.5))
    tx_group1.name = "!!group_name"
    tf_group1 = tx_group1.text_frame
    p_g1 = tf_group1.paragraphs[0]
    p_g1.text = "Presented by APEX TEAM  |  July 2026"
    p_g1.alignment = PP_ALIGN.CENTER
    p_g1.font.size = Pt(14)
    p_g1.font.color.rgb = GRAY

    # Stage both card containers off-screen so Morph can fly them in
    create_card_shape(slide, LEFT_MARGIN, OFF_SCREEN_TOP, CARD_WIDTH, CARD_HEIGHT,
                      border_color=CRIMSON, name="!!kpi_card_left")
    create_card_shape(slide, RIGHT_MARGIN, OFF_SCREEN_TOP, CARD_WIDTH, CARD_HEIGHT,
                      border_color=SILVER, name="!!kpi_card_right")

def build_content_pair(content):
    """For one logical content slide, emit TWO physical slides:
       Sub-slide A  →  LEFT card visible, RIGHT card off-screen
       Sub-slide B  →  Both cards visible  (RIGHT morphs into position)
    """
    # ── SUB-SLIDE A: Left card enters, right card stays off-screen ──────────
    slide_a = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide_a)
    set_morph_transition(slide_a, speed='slow')
    add_header(slide_a, content["title"], content["cat"])

    # LEFT card — in final position
    create_card_shape(slide_a, LEFT_MARGIN, CARD_TOP, CARD_WIDTH, CARD_HEIGHT,
                      border_color=content["left_color"], name="!!kpi_card_left")
    create_card_text(slide_a, LEFT_MARGIN, CARD_TOP, CARD_WIDTH, CARD_HEIGHT,
                     content["left_title"], content["left_bullets"],
                     name="!!kpi_text_left", title_color=CRIMSON,
                     text_color=content["left_text_color"])

    # RIGHT card — staged off-screen (will morph up on sub-slide B)
    create_card_shape(slide_a, RIGHT_MARGIN, OFF_SCREEN_TOP, CARD_WIDTH, CARD_HEIGHT,
                      border_color=content["right_color"], name="!!kpi_card_right")
    create_card_text(slide_a, RIGHT_MARGIN, OFF_SCREEN_TOP, CARD_WIDTH, CARD_HEIGHT,
                     content["right_title"], content["right_bullets"],
                     name="!!kpi_text_right", title_color=WHITE,
                     text_color=content["right_text_color"])

    # ── SUB-SLIDE B: Right card morphs into position ────────────────────────
    slide_b = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide_b)
    set_morph_transition(slide_b, speed='slow')
    add_header(slide_b, content["title"], content["cat"])

    # LEFT card — stays in place (same position = no visible morph movement)
    create_card_shape(slide_b, LEFT_MARGIN, CARD_TOP, CARD_WIDTH, CARD_HEIGHT,
                      border_color=content["left_color"], name="!!kpi_card_left")
    create_card_text(slide_b, LEFT_MARGIN, CARD_TOP, CARD_WIDTH, CARD_HEIGHT,
                     content["left_title"], content["left_bullets"],
                     name="!!kpi_text_left", title_color=CRIMSON,
                     text_color=content["left_text_color"])

    # RIGHT card — now in final position (Morph will animate it sliding up)
    create_card_shape(slide_b, RIGHT_MARGIN, CARD_TOP, CARD_WIDTH, CARD_HEIGHT,
                      border_color=content["right_color"], name="!!kpi_card_right")
    create_card_text(slide_b, RIGHT_MARGIN, CARD_TOP, CARD_WIDTH, CARD_HEIGHT,
                     content["right_title"], content["right_bullets"],
                     name="!!kpi_text_right", title_color=WHITE,
                     text_color=content["right_text_color"])

# ── Build the full deck ─────────────────────────────────────────────────────
build_title_slide()

for content in slides_content[2:]:  # Skip the None padding entries
    build_content_pair(content)

# Save
prs.save(str(OUTPUT_PATH))
print(f"SUCCESS: Premium executive presentation generated at {OUTPUT_PATH}")
print(f"Total slides: {len(prs.slides)} (7 logical slides × 2 sub-slides each, minus title = 13 physical slides)")

